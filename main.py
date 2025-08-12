import tkinter as tk
from tkinter import Label, Entry, Button, Text, Scrollbar, END, Toplevel, messagebox, filedialog
from tkinter import simpledialog
import pyttsx3
import speech_recognition as sr
import time
import requests
import threading
from bs4 import BeautifulSoup
import webbrowser
from datetime import datetime
from PIL import Image, ImageTk
from io import BytesIO
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE 
import os
import re
import subprocess
from tkinter import ttk
from tkinter import PhotoImage
import shutil
import platform
import cv2
import random

# Try to import ollama and handle if not available
try:
    import ollama
except ImportError:
    messagebox.showwarning("Missing Library", "Ollama is not installed. Some AI features may not work. Install with: pip install ollama")
    ollama = None

# Try to import yt_dlp
try:
    import yt_dlp
except ImportError:
    messagebox.showwarning("Missing Library", "yt-dlp is not installed. YouTube downloader won't work. Install with: pip install yt-dlp")
    yt_dlp = None

# Windows-specific imports
if platform.system() == "Windows":
    try:
        import win32gui
        import psutil
        import win32com.client as win32
    except ImportError:
        messagebox.showwarning("Missing Libraries", "For advanced Windows features, please install: pip install pywin32 psutil")
        win32gui = None
        psutil = None
        win32 = None

# Document processing imports
try:
    from docx import Document
    from docx.shared import Inches as DocxInches
except ImportError:
    messagebox.showwarning("Missing Library", "To handle Word documents, please install python-docx: pip install python-docx")
    Document = None

try:
    import pdfplumber
except ImportError:
    messagebox.showwarning("Missing Library", "To handle PDF text extraction, please install pdfplumber: pip install pdfplumber")
    pdfplumber = None

try:
    from pdf2docx import Converter
except ImportError:
    messagebox.showwarning("Missing Library", "For high-fidelity PDF to Word conversion, please install pdf2docx: pip install pdf2docx")
    Converter = None    

# Initialize text-to-speech engine
engine = pyttsx3.init()

# Global variables
history = []
user_name = None
previous_responses = set()
typing_interrupted = False
paused_output = False
text_entry = None
resume_index = 0
uploaded_files = []
last_code_generated = None
last_uploaded_file_context = None
last_converted_file_path = None

camera_capture = None
camera_window = None

side_panel_frame = None
is_side_panel_open = False

# --- NEW: Wake Word and State Management Variables ---
timeout_thread = None  # To hold the inactivity timer thread
is_active = False  # To track if VERA is in an active listening session
# --- End of New Variables ---


def change_voice(voice_name=None):
    """Change the TTS voice"""
    voices = engine.getProperty('voices')
    if voices:
        for voice in voices:
            if voice_name and voice_name in voice.name:
                engine.setProperty('voice', voice.id)
                print(f"Voice changed to {voice.name}")
                return
        print("Voice not found, using default voice.")
    else:
        print("No voices available, using default.")

# Try to set a specific voice
change_voice("Microsoft Zira Desktop")

def speak(text):
    """Convert text to speech"""
    try:
        engine.say(text)
        engine.runAndWait()
    except Exception as e:
        print(f"TTS Error: {e}")

def stop_all():
    """Stop/pause current output"""
    global typing_interrupted, paused_output, resume_index, last_response
    if not paused_output:
        engine.stop()
        typing_interrupted = True
        paused_output = True
        resume_index = chat_display.index(END).split(".")[0]
        chat_display.insert(END, "\n[System] Output paused by user.\n", 'bot')
    else:
        typing_interrupted = False
        paused_output = False
        chat_display.insert(END, "\n[System] Resuming output...\n", 'bot')
        if 'last_response' in globals() and last_response:
            threading.Thread(target=lambda: resumed_output(last_response)).start()
    chat_display.see(END)

def new_chat():
    """Start a new chat session"""
    chat_display.delete(1.0, END)
    global typing_interrupted, paused_output, uploaded_files, last_uploaded_file_context, last_converted_file_path, last_code_generated
    typing_interrupted = False
    paused_output = False
    uploaded_files = []
    last_code_generated = None
    last_uploaded_file_context = None
    last_converted_file_path = None
    greet_user()
    if is_side_panel_open:
        toggle_side_panel()

def resumed_output(text):
    """Resume paused output"""
    global typing_interrupted
    if text and not typing_interrupted:
        type_out_text(text)
        if not typing_interrupted:
            speak(text)

def greet_user():
    """Initial greeting message (Modified for wake word)"""
    greeting_message = "Hello sir! Say 'Hey Vera' to get started."
    chat_display.insert(END, f"[{datetime.now().strftime('%H:%M')}] VERA: {greeting_message}\n", 'bot')
    # speak(greeting_message) # Don't speak on startup, wait for wake word

def ollama_search(query, is_create_ppt=False):
    """Search using Ollama AI model"""
    if not ollama:
        return "Ollama is not available. Please install it to use AI features."
    
    try:
        response = ollama.chat(
            model='llama3.2',
            messages=[{'role': 'user', 'content': query}]
        )

        if 'message' in response and 'content' in response['message']:
            result = response['message']['content']

            if result in previous_responses:
                return "Unfortunately, no new content is available from Ollama. Please ask something else."

            if is_create_ppt:
                paragraphs = [p.strip() for p in result.split("\n\n") if p.strip()][:10]
                response_text = "\n\n".join(paragraphs)
            else:
                response_text = result

            previous_responses.add(result)
            return response_text
        else:
            return "No response content received."

    except Exception as e:
        return f"An error occurred while querying the Ollama API: {str(e)}"

def listen_command():
    """Listen for voice commands (Modified with timeout)"""
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        chat_display.insert(END, "\nVERA: Listening...", 'bot')
        chat_display.see(END)
        recognizer.adjust_for_ambient_noise(source, duration=0.5)
        try:
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=10)
            command = recognizer.recognize_google(audio)
            return command
        except sr.WaitTimeoutError:
            return "Sorry, I didn't hear anything."
        except sr.UnknownValueError:
            return "Sorry, I did not understand that."
        except sr.RequestError as e:
            return f"Could not request results; {e}"

def save_history():
    """Save search history to file"""
    try:
        with open("search_history.json", "w") as f:
            json.dump(history, f)
    except Exception as e:
        print(f"Error saving history: {e}")

def load_history():
    """Load search history from file"""
    global history
    try:
        with open("search_history.json", "r") as f:
            history = json.load(f)
    except FileNotFoundError:
        history = []
    except Exception as e:
        print(f"Error loading history: {e}")
        history = []

def delete_history():
    """Delete all search history"""
    global history
    history.clear()
    save_history()
    messagebox.showinfo("History Deleted", "All search history has been deleted.")
    show_history()

def on_history_click(event):
    """Handle clicks on history items"""
    try:
        index = history_display.index("@%s,%s" % (event.x, event.y))
        line_number = int(index.split(".")[0]) - 1
        if 0 <= line_number < len(history):
            clicked_command = history[line_number]
            execute_command(clicked_command)
    except Exception as e:
        print(f"Error handling history click: {e}")

def show_history():
    """Show search history window"""
    history_window = Toplevel(root)
    history_window.title("Search History")
    history_window.geometry("400x400")

    global history_display
    history_display = Text(history_window, wrap=tk.WORD, bg='white', fg='black', font=("Helvetica", 12))
    history_display.pack(fill=tk.BOTH, expand=True)

    scrollbar = Scrollbar(history_window, command=history_display.yview)
    scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

    history_display.config(yscrollcommand=scrollbar.set)

    for entry in history:
        history_display.insert(END, f"{entry}\n")

    history_display.insert(END, "End of history.\n")
    history_display.bind("<Button-1>", on_history_click)

    delete_button = Button(history_window, text="Delete History", command=delete_history, bg='red', fg='white')
    delete_button.pack(pady=10)

def type_out_text(text, delay=0.005):
    """Type out text with animation effect"""
    global typing_interrupted
    typing_interrupted = False
    for char in text:
        if typing_interrupted:
            break
        chat_display.insert(END, char)
        chat_display.see(END)
        chat_display.update_idletasks()
        time.sleep(delay)

def delete_item(path):
    """Delete a file or empty folder"""
    if not os.path.exists(path):
        return f"Error: '{path}' does not exist."
    try:
        if os.path.isfile(path):
            os.remove(path)
            return f"Successfully deleted file: '{path}'."
        elif os.path.isdir(path):
            if not os.listdir(path):
                os.rmdir(path)
                return f"Successfully deleted empty folder: '{path}'."
            else:
                return f"Error: Folder '{path}' is not empty. To delete non-empty folders, use 'delete folder forcefully [path]'."
    except OSError as e:
        return f"Error deleting '{path}': {e}"

def delete_folder_forcefully(path):
    """Forcefully delete a folder and its contents"""
    if not os.path.exists(path):
        return f"Error: '{path}' does not exist."
    if not os.path.isdir(path):
        return f"Error: '{path}' is not a folder."
    try:
        shutil.rmtree(path)
        return f"Successfully deleted folder and its contents: '{path}'."
    except OSError as e:
        return f"Error forcefully deleting '{path}': {e}"

def copy_item(source, destination):
    """Copy a file or folder"""
    if not os.path.exists(source):
        return f"Error: Source '{source}' does not exist."
    try:
        if os.path.isfile(source):
            shutil.copy2(source, destination)
            return f"Successfully copied file from '{source}' to '{destination}'."
        elif os.path.isdir(source):
            if os.path.exists(destination):
                return f"Error: Destination folder '{destination}' already exists. Please provide a new destination or delete the existing one."
            shutil.copytree(source, destination)
            return f"Successfully copied folder from '{source}' to '{destination}'."
    except shutil.Error as e:
        return f"Error copying '{source}' to '{destination}': {e}"
    except OSError as e:
        return f"Error copying '{source}' to '{destination}': {e}"

def move_item(source, destination):
    """Move a file or folder"""
    if not os.path.exists(source):
        return f"Error: Source '{source}' does not exist."
    try:
        shutil.move(source, destination)
        return f"Successfully moved '{source}' to '{destination}'."
    except shutil.Error as e:
        return f"Error moving '{source}' to '{destination}': {e}"
    except OSError as e:
        return f"Error moving '{source}' to '{destination}': {e}"

def create_folder(path):
    """Create a new folder"""
    if os.path.exists(path):
        return f"Error: Folder '{path}' already exists."
    try:
        os.makedirs(path)
        return f"Successfully created folder: '{path}'."
    except OSError as e:
        return f"Error creating folder '{path}': {e}"

def create_file(path, content=""):
    """Create a new file"""
    if os.path.exists(path):
        return f"Error: File '{path}' already exists."
    try:
        with open(path, 'w') as f:
            f.write(content)
        return f"Successfully created file: '{path}'."
    except OSError as e:
        return f"Error creating file '{path}': {e}"

def rename_item(old_path):
    """Rename a file or folder"""
    if not os.path.exists(old_path):
        return f"Error: '{old_path}' does not exist."

    new_name = simpledialog.askstring("Rename Item", f"Enter new name for '{os.path.basename(old_path)}':")

    if new_name:
        new_path = os.path.join(os.path.dirname(old_path), new_name)
        try:
            os.rename(old_path, new_path)
            return f"Successfully renamed '{os.path.basename(old_path)}' to '{new_name}'."
        except OSError as e:
            return f"Error renaming '{old_path}' to '{new_path}': {e}"
    else:
        return "Rename operation cancelled."

def open_application_by_name(app_name):
    """Open an application by name across different operating systems"""
    app_name_lower = app_name.lower()
    system = platform.system().lower()

    if system == "windows":
        # Common Windows application paths and executable names
        common_paths = [
            os.path.join(os.environ.get("ProgramFiles", ""), app_name),
            os.path.join(os.environ.get("ProgramFiles(x86)", ""), app_name),
            os.path.join(os.environ.get("LOCALAPPDATA", ""), app_name),
            os.path.join(os.environ.get("APPDATA", ""), app_name),
            os.path.join(os.environ.get("ProgramFiles", ""), f"{app_name} (x86)"),
            os.path.join(os.environ.get("ProgramFiles", ""), f"Google\\Chrome\\Application"),
            os.path.join(os.environ.get("ProgramFiles(x86)", ""), f"Google\\Chrome\\Application"),
            os.path.join(os.environ.get("ProgramFiles", ""), f"Microsoft\\Edge\\Application"),
            os.path.join(os.environ.get("ProgramFiles(x86)", ""), f"Microsoft\\Edge\\Application"),
            os.path.join(os.environ.get("ProgramFiles", ""), f"Mozilla Firefox"),
            os.path.join(os.environ.get("ProgramFiles(x86)", ""), f"Mozilla Firefox"),
            os.path.join(os.environ.get("ProgramFiles", ""), f"Microsoft Office\\root\\Office16"),
            os.path.join(os.environ.get("ProgramFiles(x86)", ""), f"Microsoft Office\\root\\Office16"),
        ]

        executables = {
            "chrome": ["chrome.exe"],
            "edge": ["msedge.exe"],
            "firefox": ["firefox.exe"],
            "notepad": ["notepad.exe"],
            "word": ["winword.exe"],
            "excel": ["excel.exe"],
            "powerpoint": ["powerpnt.exe"],
            "calculator": ["calc.exe"],
            "vlc": ["vlc.exe"],
            "spotify": ["spotify.exe"],
            "code": ["code.exe"],
            "discord": ["discord.exe"],
            "steam": ["steam.exe"],
        }

        # Add .exe suffix if not present
        if not app_name_lower.endswith(".exe"):
            app_name_lower_exe = app_name_lower + ".exe"
        else:
            app_name_lower_exe = app_name_lower

        # Check for direct executable match in common paths
        for path in common_paths:
            full_path = os.path.join(path, app_name_lower_exe)
            if os.path.exists(full_path):
                try:
                    subprocess.Popen(full_path)
                    return f"Successfully opened {app_name}."
                except Exception as e:
                    return f"Error opening {app_name} from {full_path}: {e}"

        # Check for known executables
        if app_name_lower in executables:
            for exe in executables[app_name_lower]:
                try:
                    subprocess.Popen(["start", exe], shell=True)
                    return f"Successfully opened {app_name}."
                except Exception as e:
                    print(f"Attempt to open {exe} directly failed: {e}")
                
                # Check in common paths for specific executables
                for path in common_paths:
                    full_path = os.path.join(path, exe)
                    if os.path.exists(full_path):
                        try:
                            subprocess.Popen(full_path)
                            return f"Successfully opened {app_name}."
                        except Exception as e:
                            return f"Error opening {app_name} from {full_path}: {e}"
        
        # Fallback: Try opening directly using os.startfile
        try:
            os.startfile(app_name)
            return f"Attempted to open {app_name} using system default. If it's installed, it should open."
        except Exception as e:
            return f"Could not find or open '{app_name}' on Windows. Please ensure it's installed and in your system's PATH, or provide the full path."

    elif system == "darwin":  # macOS
        common_paths = [
            "/Applications",
            os.path.expanduser("~/Applications"),
        ]
        
        # Add .app suffix if not present
        if not app_name_lower.endswith(".app"):
            app_name_to_open = app_name + ".app"
        else:
            app_name_to_open = app_name

        for path in common_paths:
            full_path = os.path.join(path, app_name_to_open)
            if os.path.exists(full_path):
                try:
                    subprocess.Popen(["open", full_path])
                    return f"Successfully opened {app_name}."
                except Exception as e:
                    return f"Error opening {app_name} from {full_path}: {e}"
        
        # Fallback: try opening directly with 'open' command
        try:
            subprocess.Popen(["open", "-a", app_name])
            return f"Attempted to open {app_name} using 'open -a' command. If it's installed, it should open."
        except Exception as e:
            return f"Could not find or open '{app_name}' on macOS."

    elif system == "linux":
        # Common Linux application commands
        executables = {
            "chrome": ["google-chrome", "google-chrome-stable"],
            "firefox": ["firefox"],
            "code": ["code"],
            "vlc": ["vlc"],
            "gnome-calculator": ["gnome-calculator"],
            "gedit": ["gedit"],
            "libreoffice writer": ["libreoffice", "--writer"],
            "libreoffice calc": ["libreoffice", "--calc"],
            "libreoffice impress": ["libreoffice", "--impress"],
        }

        if app_name_lower in executables:
            for cmd in executables[app_name_lower]:
                try:
                    subprocess.Popen([cmd])
                    return f"Successfully opened {app_name}."
                except FileNotFoundError:
                    continue
                except Exception as e:
                    return f"Error opening {app_name} with command '{cmd}': {e}"
        
        # Fallback: try opening directly
        try:
            subprocess.Popen([app_name_lower])
            return f"Attempted to open {app_name} directly. If it's in your PATH, it should open."
        except FileNotFoundError:
            return f"Could not find or open '{app_name}' on Linux. Please ensure it's installed and in your system's PATH."
        except Exception as e:
            return f"Error opening {app_name}: {e}"

    return f"Could not find or open '{app_name}' on your operating system ({system})."

def execute_command(command=None):
    """Main command execution function (Modified for wake word and inactivity timer)"""
    global typing_interrupted, last_response, last_uploaded_file_context, last_converted_file_path, last_code_generated, is_active, timeout_thread
    
    # A new command is issued, so VERA is active. Cancel any scheduled sleep.
    is_active = True
    if timeout_thread and timeout_thread.is_alive():
        timeout_thread.cancel()
        print("Inactivity timer cancelled by new command.")

    typing_interrupted = False

    if not command:
        command = listen_command()

    # Handle cases where listening timed out or was misunderstood
    if "sorry, i didn't hear anything" in command.lower():
        chat_display.delete("end-2l", END)  # Remove "Listening..."
        speak(command)
        # Schedule going to sleep
        timeout_thread = threading.Timer(10.0, lambda: root.after(0, go_to_sleep))
        timeout_thread.start()
        return
    elif "sorry, i did not understand that" in command.lower():
        chat_display.insert(END, f"\n[{datetime.now().strftime('%H:%M')}] VERA: {command}\n", 'bot')
        speak(command)
        # Schedule going to sleep
        timeout_thread = threading.Timer(10.0, lambda: root.after(0, go_to_sleep))
        timeout_thread.start()
        return
    elif "terminate yourself" in command.lower():
        speak("Terminating myself. Goodbye sir.")
        root.after(2000, on_closing) # Safely close the app
        return

    timestamp = datetime.now().strftime('%H:%M')
    chat_display.delete("end-2l", END)
    chat_display.insert(END, f"\n[{timestamp}] You: {command}\n", 'user')
    chat_display.insert(END, f"[{timestamp}] VERA: Processing...\n", 'bot')
    chat_display.see(END)

    loading_label, spinner = show_loading_indicator()

    def perform_search():
        global last_response, last_uploaded_file_context, last_converted_file_path, last_code_generated, timeout_thread
        try:
            speak("Please wait, I am processing your request...")
            command_lower = command.lower()

            # Handle contextual commands for uploaded files
            if last_uploaded_file_context:
                path = last_uploaded_file_context.get('path')
                file_type = last_uploaded_file_context.get('type')
                cmd = command.lower()
                
                context_command_handled = True
                
                if "explain" in cmd and file_type == 'ppt':
                    result = explain_ppt(path)
                elif "convert to word" in cmd and file_type == 'ppt':
                    new_path = convert_ppt_to_word(path)
                    if new_path and os.path.exists(new_path) and not new_path.startswith("Error"):
                        last_converted_file_path = new_path
                        result = "Successfully converted PPT to Word. To save it to your desktop, just say 'save it'."
                    else:
                        result = new_path or "Conversion failed"
                elif "convert to word" in cmd and file_type == 'pdf':
                    new_path = convert_pdf_to_word(path)
                    if new_path and os.path.exists(new_path) and not new_path.startswith("Error"):
                        last_converted_file_path = new_path
                        result = "Successfully converted PDF to Word. To save it to your desktop, just say 'save it'."
                    else:
                        result = new_path or "Conversion failed"
                elif "convert to pdf" in cmd and (file_type == 'ppt' or file_type == 'word'):
                    if file_type == 'ppt':
                        new_path = convert_ppt_to_pdf(path)
                    else:
                        new_path = convert_word_to_pdf(path)
                    
                    if new_path and os.path.exists(new_path) and not new_path.startswith("Error"):
                        last_converted_file_path = new_path
                        result = "Successfully converted to PDF. To save it to your desktop, just say 'save it'."
                    else:
                        result = new_path or "Conversion failed"
                else:
                    context_command_handled = False
                
                if context_command_handled:
                    last_uploaded_file_context = None
                else:
                    result = ollama_search(command)

            elif "save it" in command.lower() and last_converted_file_path:
                path_to_save = last_converted_file_path
                last_converted_file_path = None
                result = save_to_desktop(path_to_save)

            # Handle coding requests
            elif any(phrase in command.lower() for phrase in ['code for', 'program for', 'script for', 'write a program', 'write a script', 'generate code']):
                code_prompt = f"Write a Python script for the following request. Provide only the raw code, without any surrounding text, explanations, or markdown formatting like ```python. The request is: '{command}'"
                
                generated_code = ollama_search(code_prompt)

                if "Unfortunately" in generated_code or not generated_code.strip():
                    result = "I'm sorry, I couldn't generate the code for that request. Please try rephrasing it."
                else:
                    last_code_generated = generated_code
                    
                    upload_dir = os.path.join(os.path.expanduser("~"), "VERA_uploads")
                    os.makedirs(upload_dir, exist_ok=True)
                    
                    code_file_path = os.path.join(upload_dir, "vera_generated_code.txt")
                    
                    with open(code_file_path, "w", encoding='utf-8') as f:
                        f.write(generated_code)
                        
                    subprocess.Popen(['notepad.exe', code_file_path])
                    result = "I have written the code for you. If you want to copy it, just say 'copy it'."

            # Handle copy code command
            elif "copy it" in command.lower():
                if last_code_generated:
                    root.clipboard_clear()
                    root.clipboard_append(last_code_generated)
                    result = "The code has been copied to your clipboard."
                else:
                    result = "There is no code to copy. Please ask me to write some code first."

            # Handle math expressions
            elif re.match(r'^[\d\+\-\*/\%\(\)\.\s]+$', command):
                try:
                    result = str(eval(command))
                except Exception as e:
                    result = f"Error evaluating the expression: {str(e)}"

            # Handle time and date requests
            elif "current time" in command_lower:
                current_time = datetime.now().strftime("%H:%M:%S")
                result = f"The current time is {current_time}."
            elif "current date" in command_lower:
                current_date = datetime.now().strftime("%Y-%m-%d")
                result = f"Today's date is {current_date}."

            # Handle self-introduction
            elif command_lower in ["introduce yourself", "who are you", "tell me about yourself"]:
                result = """I am Vera, a Virtual Expert for Responsive Assistance designed to help you with various tasks.
I can assist you with things like:
1. Search for things online or open websites by voice or text commands.
2. Open applications like Chrome or Edge, or even play videos on YouTube.
3. Create PowerPoint presentations based on your queries.
4. Keep track of your search history, and you can view or delete it whenever you like.
5. Process voice commands using a microphone and convert them into text.
6. Upload files and images for me to analyze or work with.
7. Open your camera and click photos if you ask.
8. Manage your files and folders, like creating, deleting, copying, moving, or renaming them.

Let me know what you'd like to do for you today!"""

            # Handle history command
            elif command_lower == "open your history":
                result = open_history_command()

            # Handle browser search commands
            elif command_lower.startswith("open edge and search for "):
                query = command[len("open edge and search for "):].strip()
                result = open_edge_search(query)
            elif command_lower.startswith("open chrome and search for "):
                query = command[len("open chrome and search for "):].strip()
                result = open_chrome_search(query)

            # Handle YouTube commands
            elif command_lower.startswith("open youtube.com and play "):
                video_name = command[len("open youtube.com and play "):].strip()
                # Corrected URL: Removed Markdown link formatting
                youtube_url = f"[https://www.youtube.com/results?search_query=](https://www.youtube.com/results?search_query=){video_name.replace(' ', '+')}"
                webbrowser.open(youtube_url)
                result = f"Opening YouTube and searching for '{video_name}'"

            # Handle camera commands
            elif command_lower.startswith("open camera"):
                result = open_camera()
            elif command_lower == "click photo" or command_lower == "click photos":
                result = click_photo()

            # Handle open commands
            elif command_lower.startswith("open "):
                target = command[len("open "):].strip()

                # Check if it's a URL
                if re.search(r'^(https?://|www\.)|(\.com|\.org|\.net|\.edu|\.gov)$', target, re.IGNORECASE):
                    url_to_open = target if target.startswith("http") else "http://" + target
                    webbrowser.open(url_to_open)
                    result = f"Opening {url_to_open}"
                else:
                    # Try to open as an application first
                    app_open_result = open_application_by_name(target)
                    if "Successfully opened" in app_open_result or "Attempted to open" in app_open_result:
                        result = app_open_result
                    else:
                        # Try to open as a file/folder
                        speak(f"Attempting to open '{target}' as a file or folder...")
                        found_path = search_for_item_utility(target)
                        if found_path:
                            try:
                                os.startfile(found_path)
                                result = f"Successfully opened: {found_path}"
                            except Exception as e:
                                result = f"Found '{found_path}', but failed to open it: {str(e)}"
                        else:
                            result = f"Could not find '{target}' as an application, file, or folder on your computer."

            # Handle close commands
            elif command_lower.startswith("close "):
                item_to_close = command[len("close "):].strip()
                result = close_application(item_to_close)

            # Handle PowerPoint creation
            elif command_lower.startswith("create ppt"):
                result = create_ppt(command)

            # Handle file analysis
            elif any(keyword in command_lower for keyword in ["analyze file", "process file", "check file", "examine file", "explain ppt"]):
                if uploaded_files:
                    last_file = uploaded_files[-1]
                    if last_file.lower().endswith('.pptx'):
                        result = explain_ppt(last_file)
                    else:
                        result = analyze_uploaded_files()
                else:
                    result = "No files have been uploaded yet. Please use the 'Add File' button to upload a file first."

            # Handle file operations
            elif command_lower.startswith("delete file "):
                file_name = command[len("delete file "):].strip()
                found_path = search_for_item_utility(file_name)
                if found_path:
                    result = delete_item(found_path)
                else:
                    result = f"Could not find file '{file_name}' to delete."

            elif command_lower.startswith("delete folder forcefully "):
                folder_name = command[len("delete folder forcefully "):].strip()
                found_path = search_for_item_utility(folder_name)
                if found_path:
                    result = delete_folder_forcefully(found_path)
                else:
                    result = f"Could not find folder '{folder_name}' to forcefully delete."

            elif command_lower.startswith("delete folder "):
                folder_name = command[len("delete folder "):].strip()
                found_path = search_for_item_utility(folder_name)
                if found_path:
                    result = delete_item(found_path)
                else:
                    result = f"Could not find folder '{folder_name}' to delete."

            elif command_lower.startswith("copy file ") and " to " in command_lower:
                parts = command_lower[len("copy file "):].split(" to ", 1)
                if len(parts) == 2:
                    source_name = parts[0].strip()
                    destination_path = parts[1].strip()
                    found_source_path = search_for_item_utility(source_name)
                    if found_source_path:
                        result = copy_item(found_source_path, destination_path)
                    else:
                        result = f"Could not find source file '{source_name}' to copy."
                else:
                    result = "Invalid 'copy file' command format. Use: 'copy file [source_name] to [destination_path]'."

            elif command_lower.startswith("copy folder ") and " to " in command_lower:
                parts = command_lower[len("copy folder "):].split(" to ", 1)
                if len(parts) == 2:
                    source_name = parts[0].strip()
                    destination_path = parts[1].strip()
                    found_source_path = search_for_item_utility(source_name)
                    if found_source_path:
                        result = copy_item(found_source_path, destination_path)
                    else:
                        result = f"Could not find source folder '{source_name}' to copy."
                else:
                    result = "Invalid 'copy folder' command format. Use: 'copy folder [source_name] to [destination_path]'."

            elif command_lower.startswith("move file ") and " to " in command_lower:
                parts = command_lower[len("move file "):].split(" to ", 1)
                if len(parts) == 2:
                    source_name = parts[0].strip()
                    destination_path = parts[1].strip()
                    found_source_path = search_for_item_utility(source_name)
                    if found_source_path:
                        result = move_item(found_source_path, destination_path)
                    else:
                        result = f"Could not find source file '{source_name}' to move."
                else:
                    result = "Invalid 'move file' command format. Use: 'move file [source_name] to [destination_path]'."

            elif command_lower.startswith("move folder ") and " to " in command_lower:
                parts = command_lower[len("move folder "):].split(" to ", 1)
                if len(parts) == 2:
                    source_name = parts[0].strip()
                    destination_path = parts[1].strip()
                    found_source_path = search_for_item_utility(source_name)
                    if found_source_path:
                        result = move_item(found_source_path, destination_path)
                    else:
                        result = f"Could not find source folder '{source_name}' to move."
                else:
                    result = "Invalid 'move folder' command format. Use: 'move folder [source_name] to [destination_path]'."

            elif command_lower.startswith("create folder "):
                folder_path = command[len("create folder "):].strip()
                result = create_folder(folder_path)

            elif command_lower.startswith("create file "):
                file_path = command[len("create file "):].strip()
                result = create_file(file_path)

            elif command_lower.startswith("rename file ") or command_lower.startswith("rename folder "):
                item_name = command_lower.replace("rename file ", "").replace("rename folder ", "").strip()
                found_path = search_for_item_utility(item_name)
                if found_path:
                    result = rename_item(found_path)
                else:
                    result = f"Could not find '{item_name}' to rename."

            # Default: use Ollama search
            else:
                result = ollama_search(command)

            last_response = result

            chat_display.delete("end-2l", END)
            chat_display.insert(END, f"\n[{timestamp}] VERA:", 'bot')
            chat_display.see(END)

            if not typing_interrupted:
                type_out_text(result)
                if not typing_interrupted:
                    speak(result)
                    speak("Task completed successfully.")

        except Exception as e:
            error_msg = f"An error occurred: {str(e)}"
            chat_display.insert(END, f"\n[{timestamp}] VERA: {error_msg}\n", 'bot')
            speak(error_msg)
        finally:
            hide_loading_indicator(loading_label, spinner)
            # After a task is done, start the 10-second timer to go to sleep.
            if "terminate yourself" not in command.lower():
                print("Scheduling go_to_sleep in 10 seconds.")
                timeout_thread = threading.Timer(10.0, lambda: root.after(0, go_to_sleep))
                timeout_thread.start()

    threading.Thread(target=perform_search).start()

    history.append(command)
    save_history()

# --- NEW: Wake Word and State Change Functions ---
def go_to_sleep():
    """Puts VERA into the sleeping state, waiting for the wake word."""
    global is_active
    is_active = False
    message = "\n[System] Inactivity detected. Going back to sleep. Say 'Hey Vera' to wake me up.\n"
    chat_display.insert(END, message, 'bot')
    chat_display.see(END)
    speak("Going back to sleep.")

def activate_and_listen():
    """Activates VERA to listen for a single command."""
    speak("Yes sir, how can I help you?")
    # Use a thread to call execute_command so the GUI doesn't freeze while listening
    threading.Thread(target=execute_command).start()

def listen_for_wake_word():
    """
    Runs in a background thread, continuously listening for the wake word 'Hey Vera'.
    """
    global is_active, timeout_thread
    recognizer = sr.Recognizer()
    while True:  # Loop forever
        if not is_active:
            with sr.Microphone() as source:
                try:
                    print("Listening for wake word...")
                    recognizer.adjust_for_ambient_noise(source, duration=0.5)
                    audio = recognizer.listen(source, phrase_time_limit=4)
                    command = recognizer.recognize_google(audio).lower()
                    print(f"Heard: {command}")
                    if "hey vera" in command or "hey veera" in command:
                        is_active = True
                        # Cancel any lingering timer just in case
                        if timeout_thread and timeout_thread.is_alive():
                            timeout_thread.cancel()
                        # Use root.after to safely call GUI-related functions from this thread
                        root.after(0, activate_and_listen)
                except sr.UnknownValueError:
                    # Silence or unrecognized speech, continue listening
                    pass
                except sr.RequestError as e:
                    print(f"Could not request results from Google; {e}")
                    time.sleep(5) # Wait before retrying
                except Exception as e:
                    print(f"An error occurred in wake word listener: {e}")
        else:
            # If VERA is active, pause this loop to save resources
            time.sleep(0.5)
# --- End of New Functions ---

def open_camera():
    """Open camera window"""
    global camera_capture, camera_window

    if camera_capture is not None and camera_capture.isOpened():
        return "Your camera is already open."

    camera_capture = cv2.VideoCapture(0)

    if not camera_capture.isOpened():
        return "Could not open camera. Please ensure it's connected and not in use by another application."

    camera_window = Toplevel(root)
    camera_window.title("V.E.R.A Camera Feed")
    camera_window.geometry("640x480")
    camera_window.protocol("WM_DELETE_WINDOW", close_camera_window)

    camera_label = Label(camera_window)
    camera_label.pack()

    def update_frame():
        if camera_capture and camera_capture.isOpened():
            ret, frame = camera_capture.read()
            if ret:
                cv2image = cv2.cvtColor(frame, cv2.COLOR_BGR2RGBA)
                img = Image.fromarray(cv2image)
                imgtk = ImageTk.PhotoImage(image=img)
                camera_label.imgtk = imgtk
                camera_label.configure(image=imgtk)
                camera_label.after(10, update_frame)
            else:
                messagebox.showwarning("Camera Error", "Failed to capture frame from camera.")
                close_camera_window()

    update_frame()
    return "Opening your camera now."

def close_camera_window():
    """Close camera window and release camera"""
    global camera_capture, camera_window
    if camera_capture:
        camera_capture.release()
        camera_capture = None
    if camera_window:
        camera_window.destroy()
        camera_window = None
    chat_display.insert(END, "\n[VERA] Camera closed.\n", 'bot')
    speak("Camera closed.")

def click_photo():
    """Capture a photo from the camera"""
    global camera_capture

    if camera_capture is None or not camera_capture.isOpened():
        return "Camera is not open. Please say 'open camera' first."

    ret, frame = camera_capture.read()
    if ret:
        timestamp_str = datetime.now().strftime("%Y%m%d_%H%M%S")
        pictures_dir = os.path.join(os.path.expanduser("~"), "Pictures")
        os.makedirs(pictures_dir, exist_ok=True)
        photo_filename = os.path.join(pictures_dir, f"VERA_Photo_{timestamp_str}.png")
        cv2.imwrite(photo_filename, frame)
        return f"Photo saved to: {photo_filename}"
    else:
        return "Failed to capture photo. Please try again."

def add_footer(slide, slide_number, total_slides, text_color):
    """Add footer to PowerPoint slide"""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(9.0)
    height = Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, width / 2, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Generated by V.E.R.A"
    p.font.size = Pt(10)
    p.font.color.rgb = text_color

    txBox = slide.shapes.add_textbox(left + width / 2, top, width / 2, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Slide {slide_number} of {total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = text_color

def explain_ppt(file_path):
    """Explain PowerPoint presentation slides"""
    try:
        if not file_path.lower().endswith('.pptx'):
            return "This function only works with .pptx files."

        chat_display.insert(END, f"\n[VERA] Starting slideshow for {os.path.basename(file_path)}...\n", 'bot')
        speak(f"Starting slideshow for the presentation {os.path.basename(file_path)}.")

        system = platform.system().lower()
        slideshow_started = False
        powerpoint_app = None
        presentation = None
        slideshow_window = None

        if system == "windows" and win32:
            try:
                powerpoint_app = win32.Dispatch("PowerPoint.Application")
                powerpoint_app.Visible = True

                presentation = powerpoint_app.Presentations.Open(file_path)
                slideshow_window = presentation.SlideShowSettings.Run()
                slideshow_started = True
                chat_display.insert(END, f"\n[VERA] PowerPoint slideshow started.\n", 'bot')
            except Exception as e:
                speak(f"Could not automatically start the slideshow using COM: {e}. Opening file normally.")
                chat_display.insert(END, f"\n[VERA] Could not start slideshow automatically: {e}\n", 'bot')
                try:
                    command = f'start powerpnt /s "{file_path}"'
                    subprocess.Popen(command, shell=True)
                    slideshow_started = True
                except Exception as sub_e:
                    chat_display.insert(END, f"\n[VERA] Fallback open failed: {sub_e}\n", 'bot')
                slideshow_started = False
        elif system == "darwin":
            subprocess.Popen(["open", file_path])
            speak("The presentation is open. Please start the slideshow manually and navigate as I speak.")
        else:
            subprocess.Popen(["libreoffice", "--show", file_path])
            speak("The presentation is open. Please start the slideshow manually and navigate as I speak.")

        if slideshow_started:
            speak("Giving the presentation a moment to load.")
            time.sleep(5)

        prs = Presentation(file_path)
        chat_display.insert(END, f"\n[VERA] Starting explanation...\n", 'bot')
        speak("Starting explanation.")

        # Save current engine properties
        original_rate = engine.getProperty('rate')
        original_volume = engine.getProperty('volume')

        # Set human-like properties for explanation
        engine.setProperty('rate', 175)
        engine.setProperty('volume', 0.9)

        for i, slide in enumerate(prs.slides):
            if typing_interrupted:
                speak("Explanation stopped by user.")
                if system == "windows" and powerpoint_app:
                    try:
                        powerpoint_app.Quit()
                    except:
                        pass
                engine.setProperty('rate', original_rate)
                engine.setProperty('volume', original_volume)
                return "Explanation stopped."

            slide_number = i + 1
            chat_display.insert(END, f"\n--- Slide {slide_number} ---\n", 'user')
            speak(f"Slide {slide_number}")

            if system == "windows" and slideshow_started and slideshow_window:
                try:
                    slideshow_view = slideshow_window.View
                    slideshow_view.GotoSlide(slide_number)
                    print(f"Navigated to slide {slide_number}")
                except Exception as nav_e:
                    print(f"Error navigating to slide {slide_number}: {nav_e}")
                    chat_display.insert(END, f"[VERA] Could not navigate to slide {slide_number}: {nav_e}\n", 'bot')
                    speak(f"I encountered an issue navigating to slide {slide_number}.")

            slide_text = []
            if slide.shapes.title:
                slide_text.append(slide.shapes.title.text)

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)

            full_slide_text = "\n".join(slide_text).strip()

            if not full_slide_text:
                chat_display.insert(END, "[Content]: No text found on this slide.\n", 'bot')
                speak("No text found on this slide.")
                time.sleep(5)
                continue

            chat_display.insert(END, f"[Content]: {full_slide_text}\n", 'bot')

            explanation_prompt = f"Explain the following content from a presentation slide in a clear and concise way: {full_slide_text}"
            explanation = ollama_search(explanation_prompt)

            chat_display.insert(END, f"[Explanation]: {explanation}\n", 'bot')
            speak(explanation)

            time.sleep(5)

        final_message = "Finished explaining the presentation."
        speak(final_message)

        if system == "windows" and powerpoint_app:
            try:
                powerpoint_app.Quit()
            except:
                pass

        # Restore original engine properties
        engine.setProperty('rate', original_rate)
        engine.setProperty('volume', original_volume)
        
        return final_message

    except Exception as e:
        error_message = f"An error occurred while explaining the PowerPoint: {str(e)}"
        speak(error_message)
        if system == "windows" and 'powerpoint_app' in locals() and powerpoint_app:
            try:
                powerpoint_app.Quit()
            except:
                pass
        return error_message

def analyze_uploaded_files():
    """Analyze uploaded files and provide information"""
    if not uploaded_files:
        return "No files are currently uploaded."

    analysis = "Here's what I found in your uploaded files:\n\n"

    for file_path in uploaded_files:
        filename = os.path.basename(file_path)
        extension = os.path.splitext(file_path)[1].lower()

        analysis += f"File: {filename}\n"
        analysis += f"Type: {extension[1:].upper()} file\n"

        try:
            file_size = os.path.getsize(file_path)
            if file_size < 1024:
                size_str = f"{file_size} bytes"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.2f} MB"

            analysis += f"Size: {size_str}\n"

            if extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                try:
                    img = Image.open(file_path)
                    analysis += f"Image dimensions: {img.width}x{img.height} pixels\n"
                    analysis += f"Image mode: {img.mode}\n"
                except Exception as e:
                    analysis += f"Could not analyze image details: {str(e)}\n"

            analysis += "\n"

        except Exception as e:
            analysis += f"Error analyzing this file: {str(e)}\n\n"

    return analysis

def upload_file():
    """Upload files to VERA"""
    global uploaded_files, last_uploaded_file_context

    file_paths = filedialog.askopenfilenames(
        title="Select Files",
        filetypes=(
            ("PowerPoint files", "*.pptx"),
            ("Image files", "*.jpg *.jpeg *.png *.gif *.bmp"),
            ("Document files", "*.pdf *.doc *.docx *.txt"),
            ("All files", "*.*")
        )
    )

    if not file_paths:
        return

    timestamp = datetime.now().strftime('%H:%M')
    upload_dir = os.path.join(os.path.expanduser("~"), "VERA_uploads")
    os.makedirs(upload_dir, exist_ok=True)

    for file_path in file_paths:
        filename = os.path.basename(file_path)
        destination = os.path.join(upload_dir, filename)
        shutil.copy2(file_path, destination)
        uploaded_files.append(destination)

        chat_display.insert(END, f"\n[{timestamp}] You: [Uploaded file: {filename}]\n", 'user')

        extension = os.path.splitext(file_path)[1].lower()
        response_message = ""
        
        if extension == '.pptx':
            last_uploaded_file_context = {'path': destination, 'type': 'ppt'}
            response_message = "PowerPoint file detected. You can ask me to 'explain the presentation', 'convert to word', or 'convert to pdf'."
        elif extension in ['.doc', '.docx']:
            last_uploaded_file_context = {'path': destination, 'type': 'word'}
            response_message = "Word document detected. You can ask me to 'convert to pdf'."
        elif extension == '.pdf':
            last_uploaded_file_context = {'path': destination, 'type': 'pdf'}
            response_message = "PDF file detected. You can ask me to 'convert to word'."
        elif extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            last_uploaded_file_context = None
            response_message = "Image file uploaded. You can ask me to analyze it."
            try:
                img = Image.open(file_path)
                img.thumbnail((300, 300))
                photo = ImageTk.PhotoImage(img)
                chat_display.insert(END, "\n")
                chat_display.image_create(END, image=photo)
                chat_display.insert(END, "\n")
                chat_display.image_refs = getattr(chat_display, 'image_refs', []) + [photo]
            except Exception as e:
                chat_display.insert(END, f"[Error displaying image preview: {str(e)}]\n")
        else:
            last_uploaded_file_context = None
            response_message = f"{filename} uploaded. You can ask me to analyze it."

        chat_display.insert(END, f"\n[{timestamp}] VERA: {response_message}\n", 'bot')
        speak("File uploaded.")

    chat_display.see(END)

def convert_ppt_to_pdf(ppt_path):
    """Convert PowerPoint to PDF"""
    if platform.system() != "Windows" or not win32:
        return "Error: PPT to PDF conversion is only supported on Windows with PowerPoint installed."
    
    powerpoint_app = None
    try:
        pdf_path = os.path.splitext(ppt_path)[0] + ".pdf"
        powerpoint_app = win32.Dispatch("PowerPoint.Application")
        presentation = powerpoint_app.Presentations.Open(os.path.abspath(ppt_path))
        presentation.SaveAs(os.path.abspath(pdf_path), 32)  # 32 is ppSaveAsPDF
        presentation.Close()
        return pdf_path
    except Exception as e:
        return f"Error during PPT to PDF conversion: {e}"
    finally:
        if powerpoint_app:
            try:
                powerpoint_app.Quit()
            except:
                pass

def convert_word_to_pdf(doc_path):
    """Convert Word document to PDF"""
    if platform.system() != "Windows" or not win32:
        return "Error: Word to PDF conversion is only supported on Windows with Word installed."
        
    word_app = None
    try:
        pdf_path = os.path.splitext(doc_path)[0] + ".pdf"
        word_app = win32.Dispatch("Word.Application")
        doc = word_app.Documents.Open(os.path.abspath(doc_path))
        doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)  # 17 is wdFormatPDF
        doc.Close()
        return pdf_path
    except Exception as e:
        return f"Error during Word to PDF conversion: {e}"
    finally:
        if word_app:
            try:
                word_app.Quit()
            except:
                pass

def convert_pdf_to_word(pdf_path):
    """Convert PDF to DOCX using high-fidelity library"""
    if not Converter:
        return "Error: The pdf2docx library is required for this conversion. Please install it."
    
    try:
        doc_path = os.path.splitext(pdf_path)[0] + ".docx"
        cv = Converter(pdf_path)
        cv.convert(doc_path, start=0, end=None)
        cv.close()
        return doc_path
    except Exception as e:
        return f"Error during high-fidelity PDF to Word conversion: {e}"

def convert_ppt_to_word(ppt_path):
    """Convert PPTX to DOCX, including text, images, and tables"""
    if not Document:
        return "Error: The python-docx library is required. Please install it."
    
    try:
        prs = Presentation(ppt_path)
        doc = Document()
        doc_path = os.path.splitext(ppt_path)[0] + ".docx"

        for i, slide in enumerate(prs.slides):
            doc.add_heading(f'--- Slide {i+1} ---', level=2)
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    doc.add_paragraph(shape.text)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_stream = BytesIO(shape.image.blob)
                    try:
                        doc.add_picture(image_stream, width=DocxInches(5.0))
                    except Exception as e:
                        print(f"Could not add image from slide {i+1}: {e}")

                if shape.has_table:
                    table = shape.table
                    doc_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                    doc_table.style = 'Table Grid'
                    for r, row in enumerate(table.rows):
                        for c, cell in enumerate(row.cells):
                            doc_table.cell(r, c).text = cell.text
            
            doc.add_page_break()

        doc.save(doc_path)
        return doc_path
    except Exception as e:
        return f"Error during enhanced PPT to Word conversion: {e}"

def save_to_desktop(source_path):
    """Save file to desktop"""
    if not source_path or not os.path.exists(source_path):
        return "Error: No file to save or source file not found."
    
    try:
        desktop_path = os.path.join(os.path.expanduser("~"), "Desktop")
        filename = os.path.basename(source_path)
        destination_path = os.path.join(desktop_path, filename)
        
        if os.path.exists(destination_path):
            base, ext = os.path.splitext(filename)
            i = 1
            while os.path.exists(destination_path):
                destination_path = os.path.join(desktop_path, f"{base}_{i}{ext}")
                i += 1
                
        shutil.move(source_path, destination_path)
        return f"File saved successfully to your desktop as {os.path.basename(destination_path)}."
    except Exception as e:
        return f"Error: Failed to save file to desktop: {e}"

def open_chrome_search(query):
    """Open Chrome and search for query"""
    try:
        search_url = f"[https://www.google.com/search?q=](https://www.google.com/search?q=){query.replace(' ', '+')}"
        subprocess.run(["start", "chrome", search_url], shell=True)
        return f"Opening Google Chrome and searching for '{query}'"
    except Exception as e:
        return f"An error occurred while opening Google Chrome: {str(e)}"

def open_edge_search(query):
    """Open Edge and search for query"""
    try:
        edge_path = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
        search_url = f"[https://www.google.com/search?q=](https://www.google.com/search?q=){query.replace(' ', '+')}"
        subprocess.run([edge_path, search_url], shell=True)
        return f"Opening Microsoft Edge and searching for '{query}'"
    except Exception as e:
        return f"An error occurred while opening Microsoft Edge: {str(e)}"

# Color palettes for PowerPoint presentations
COLOR_PALETTES = [
    {
        "PRIMARY_COLOR": RGBColor(38, 70, 83),
        "ACCENT_COLOR": RGBColor(42, 157, 143),
        "LIGHT_BACKGROUND": RGBColor(244, 247, 245),
        "DARK_BACKGROUND": RGBColor(38, 70, 83),
        "TEXT_LIGHT": RGBColor(233, 236, 239),
        "TEXT_DARK": RGBColor(0, 0, 0),
        "FOOTER_TEXT": RGBColor(100, 100, 100)
    },
    {
        "PRIMARY_COLOR": RGBColor(25, 118, 210),
        "ACCENT_COLOR": RGBColor(255, 152, 0),
        "LIGHT_BACKGROUND": RGBColor(227, 242, 253),
        "DARK_BACKGROUND": RGBColor(25, 118, 210),
        "TEXT_LIGHT": RGBColor(255, 255, 255),
        "TEXT_DARK": RGBColor(33, 33, 33),
        "FOOTER_TEXT": RGBColor(80, 80, 80)
    },
    {
        "PRIMARY_COLOR": RGBColor(67, 160, 71),
        "ACCENT_COLOR": RGBColor(123, 31, 162),
        "LIGHT_BACKGROUND": RGBColor(232, 245, 233),
        "DARK_BACKGROUND": RGBColor(67, 160, 71),
        "TEXT_LIGHT": RGBColor(255, 255, 255),
        "TEXT_DARK": RGBColor(51, 51, 51),
        "FOOTER_TEXT": RGBColor(90, 90, 90)
    },
    {
        "PRIMARY_COLOR": RGBColor(211, 47, 47),
        "ACCENT_COLOR": RGBColor(96, 125, 139),
        "LIGHT_BACKGROUND": RGBColor(250, 240, 240),
        "DARK_BACKGROUND": RGBColor(211, 47, 47),
        "TEXT_LIGHT": RGBColor(255, 255, 255),
        "TEXT_DARK": RGBColor(60, 60, 60),
        "FOOTER_TEXT": RGBColor(120, 120, 120)
    },
    {
        "PRIMARY_COLOR": RGBColor(255, 193, 7),
        "ACCENT_COLOR": RGBColor(26, 35, 126),
        "LIGHT_BACKGROUND": RGBColor(255, 248, 225),
        "DARK_BACKGROUND": RGBColor(26, 35, 126),
        "TEXT_LIGHT": RGBColor(255, 255, 255),
        "TEXT_DARK": RGBColor(40, 40, 40),
        "FOOTER_TEXT": RGBColor(110, 110, 110)
    }
]

def create_ppt(command):
    """Create PowerPoint presentation"""
    try:
        command_topic = command.lower().replace("create ppt", "").strip()
        if not command_topic:
            return "Please provide a topic for the presentation."

        llama_prompt = (
            f"Generate 8-12 concise, distinct, and informative key points or sections about '{command_topic}' "
            "suitable for a presentation. Each point should be a well-formed paragraph, separated by double newlines. "
            "Focus on providing enough detail for a slide, avoiding overly short or redundant sections."
        )
        search_results = ollama_search(llama_prompt, is_create_ppt=True)

        if "Unfortunately" in search_results or not search_results:
            return search_results if search_results else "No content generated for the presentation."

        selected_palette = random.choice(COLOR_PALETTES)
        PRIMARY_COLOR = selected_palette["PRIMARY_COLOR"]
        ACCENT_COLOR = selected_palette["ACCENT_COLOR"]
        LIGHT_BACKGROUND = selected_palette["LIGHT_BACKGROUND"]
        DARK_BACKGROUND = selected_palette["DARK_BACKGROUND"]
        TEXT_LIGHT = selected_palette["TEXT_LIGHT"]
        TEXT_DARK = selected_palette["TEXT_DARK"]
        FOOTER_TEXT = selected_palette["FOOTER_TEXT"]

        document_path = os.path.join(os.path.expanduser("~"), "Desktop")
        os.makedirs(document_path, exist_ok=True)

        ppt_file = os.path.join(document_path, f"{command_topic.replace(' ', '_')}.pptx")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        raw_responses = [p.strip() for p in search_results.split("\n\n") if p.strip()]

        grouped_content = []
        current_group = []
        for i, paragraph in enumerate(raw_responses):
            if not current_group or \
               (len(" ".join(current_group + [paragraph])) > 400 and len(current_group) > 0) or \
               (len(paragraph) > 300 and len(current_group) == 0):
                if current_group:
                    grouped_content.append(" ".join(current_group))
                current_group = [paragraph]
            else:
                current_group.append(paragraph)
        if current_group:
            grouped_content.append(" ".join(current_group))

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BACKGROUND

        title_placeholder = slide.shapes.title
        title_placeholder.text = command_topic.upper()
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(48)
        title_placeholder.text_frame.paragraphs[0].font.bold = True

        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = "Generated by V.E.R.A"
        subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)

        generic_slide_titles = [
            "Overview", "Key Aspects", "Detailed Information", "Analysis & Implications",
            "Practical Applications", "Advantages & Benefits", "Challenges & Considerations",
            "Future Perspectives", "Summary & Conclusion", "Further Insights"
        ]
        total_content_slides = len(grouped_content)
        total_slides = total_content_slides + 2

        # Content slides
        for i, content_block in enumerate(grouped_content):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = LIGHT_BACKGROUND

            first_sentence = content_block.split('.')[0].strip()
            current_slide_title = ""
            if len(first_sentence) < 60 and (first_sentence.endswith((':', '.', '!', '?')) or len(first_sentence.split()) < 10):
                current_slide_title = first_sentence.replace('.', '').replace(':', '').strip().title()
            
            if not current_slide_title:
                current_slide_title = generic_slide_titles[i % len(generic_slide_titles)]

            title_placeholder = slide.shapes.title
            title_placeholder.text = current_slide_title
            title_placeholder.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
            title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
            title_placeholder.text_frame.paragraphs[0].font.bold = True

            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()

            sentences = re.split(r'(?<=[.!?])\s+', content_block)
            sentences = [s.strip() for s in sentences if s.strip()]

            for j, sentence in enumerate(sentences):
                p = tf.add_paragraph()
                p.text = sentence
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(10)

            add_footer(slide, i + 2, total_slides, FOOTER_TEXT)

        # Thank you slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = ACCENT_COLOR

        title_placeholder = slide.shapes.title
        title_placeholder.text = "Thank You!"
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(48)
        title_placeholder.text_frame.paragraphs[0].font.bold = True

        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = "For more assistance, just ask V.E.R.A."
        subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)

        add_footer(slide, total_slides, total_slides, FOOTER_TEXT)

        prs.save(ppt_file)

        speak(f"PowerPoint presentation '{command_topic}' has been created successfully. Would you like me to open it for you?")
        open_it = messagebox.askyesno("Open Presentation", f"PowerPoint presentation '{command_topic}' has been created successfully. Would you like me to open it for you?")

        if open_it:
            os.startfile(ppt_file)
            return f"Opening presentation: '{ppt_file}'."
        else:
            return "Okay sir, I am ready for the next task."

    except Exception as e:
        return f"An error occurred while creating the PowerPoint presentation: {str(e)}"

def search_for_item_utility(item_name):
    """Search for files and folders in common locations"""
    home_dir = os.path.expanduser("~")
    search_dirs = [
        home_dir,
        os.path.join(home_dir, "Documents"),
        os.path.join(home_dir, "Downloads"),
        os.path.join(home_dir, "Desktop"),
        os.path.join(home_dir, "Pictures"),
        os.path.join(home_dir, "Music"),
        os.path.join(home_dir, "Videos"),
    ]

    # Add D: drive if it exists
    d_drive_path = "D:\\"
    if os.path.exists(d_drive_path):
        search_dirs.append(d_drive_path)
        print(f"D: drive found. Added to search paths.")

    # First, look for exact matches
    for s_dir in search_dirs:
        if not os.path.exists(s_dir):
            continue
        try:
            for item in os.listdir(s_dir):
                if item.lower() == item_name.lower():
                    full_path = os.path.join(s_dir, item)
                    return full_path
        except OSError as e:
            print(f"Could not access {s_dir}: {e}")
            continue

    # If no exact match, look for partial matches
    partial_matches = []
    for s_dir in search_dirs:
        if not os.path.exists(s_dir):
            continue
        try:
            for root, dirs, files in os.walk(s_dir):
                # Limit search depth to avoid performance issues
                if root not in search_dirs and len(root.split(os.sep)) > len(s_dir.split(os.sep)) + 2:
                    dirs[:] = []
                    continue

                for name in files + dirs:
                    if item_name.lower() in name.lower():
                        partial_matches.append(os.path.join(root, name))
        except OSError as e:
            print(f"Could not walk through {s_dir}: {e}")
            continue

    if partial_matches:
        return partial_matches[0]
    else:
        return None

def close_application(item_name):
    """Close applications, windows, and processes"""
    system = platform.system().lower()
    item_name_lower = item_name.lower()
    closed_items = []

    try:
        if system == "windows" and win32gui and psutil:
            def callback(hwnd, extra):
                if win32gui.IsWindowVisible(hwnd) and win32gui.IsWindowEnabled(hwnd):
                    title = win32gui.GetWindowText(hwnd)
                    if item_name_lower in title.lower():
                        try:
                            win32gui.PostMessage(hwnd, 0x0010, 0, 0)  # WM_CLOSE
                            closed_items.append(f"Window '{title}'")
                            print(f"Sent WM_CLOSE to: {title}")
                            time.sleep(0.5)
                            if win32gui.IsWindow(hwnd):
                                print(f"Window '{title}' did not close gracefully, attempting to force close process.")
                                tid, pid = win32gui.GetWindowThreadProcessId(hwnd)
                                try:
                                    process = psutil.Process(pid)
                                    process.terminate()
                                    closed_items[-1] += " (forced process termination)"
                                    print(f"Forced termination of process for '{title}' (PID: {pid})")
                                except psutil.NoSuchProcess:
                                    pass
                                except Exception as e:
                                    print(f"Error terminating process for '{title}' (PID: {pid}): {e}")
                        except Exception as e:
                            print(f"Error closing window '{title}': {e}")
                return True

            win32gui.EnumWindows(callback, None)

            # Close processes by name
            pids_to_kill = set()
            for proc in psutil.process_iter(['pid', 'name']):
                try:
                    if item_name_lower in proc.info['name'].lower():
                        pids_to_kill.add(proc.info['pid'])
                        if f"Process '{proc.info['name']}' (PID: {proc.info['pid']})" not in closed_items:
                            closed_items.append(f"Process '{proc.info['name']}' (PID: {proc.info['pid']})")
                            print(f"Identified process '{proc.info['name']}' (PID: {proc.info['pid']})")
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    continue

            for pid in list(pids_to_kill):
                try:
                    process = psutil.Process(pid)
                    if process.is_running():
                        process.terminate()
                        print(f"Terminated process PID {pid}")
                        # Update status in closed_items
                        for i, item in enumerate(closed_items):
                            if f"(PID: {pid})" in item:
                                closed_items[i] = item + " (terminated)"
                                break
                except psutil.NoSuchProcess:
                    for i, item in enumerate(closed_items):
                        if f"(PID: {pid})" in item:
                            closed_items[i] = item + " (already closed)"
                            break
                    print(f"Process PID {pid} already closed or not found.")
                except Exception as e:
                    print(f"Error terminating process PID {pid}: {e}")

        elif system in ["linux", "darwin"]:
            try:
                pgrep_cmd = ["pgrep", "-fi", item_name]
                pgrep_result = subprocess.run(pgrep_cmd, capture_output=True, text=True, check=True)
                pids = pgrep_result.stdout.strip().split('\n')

                if not pids or pids == ['']:
                    return f"Application/item '{item_name}' not found or is not running."

                for pid in pids:
                    if pid:
                        kill_cmd = ["kill", "-9", pid]
                        subprocess.run(kill_cmd, capture_output=True, text=True, check=False)
                        closed_items.append(f"Process (PID: {pid})")

            except subprocess.CalledProcessError as e:
                if e.returncode == 1:
                    return f"Application/item '{item_name}' not found or is not running."
                else:
                    return f"An error occurred while finding processes for {item_name}: {e.stderr.strip()}"
            except Exception as e:
                return f"An error occurred while trying to close {item_name}: {str(e)}"
        else:
            return f"Unsupported operating system: {system}"

        if closed_items:
            return f"Successfully attempted to close: {', '.join(closed_items)} related to '{item_name}'."
        else:
            return f"Could not find any open applications, files, or folders matching '{item_name}'."

    except FileNotFoundError:
        return "Error: Required system commands or Python libraries are not found on your system."
    except Exception as e:
        return f"An unexpected error occurred while trying to close {item_name}: {str(e)}"

def open_history_command():
    """Open search history"""
    speak("Opening history.")
    show_history()
    return "Opening My search history."

def show_loading_indicator():
    """Show loading animation"""
    loading_label = Label(root, text="Processing...", bg='white', font=("Helvetica", 12))
    loading_label.place(relx=0.5, rely=0.5, anchor='center')

    spinner = ttk.Progressbar(root, orient="horizontal", length=200, mode="indeterminate")
    spinner.place(relx=0.5, rely=0.6, anchor='center')
    spinner.start()
    return loading_label, spinner

def hide_loading_indicator(loading_label, spinner):
    """Hide loading animation"""
    try:
        loading_label.destroy()
        spinner.stop()
        spinner.destroy()
    except:
        pass

def on_text_command(event=None):
    """Handle text command input"""
    command = text_entry.get().strip()
    if command:
        execute_command(command)
        text_entry.delete(0, END)
    else:
        chat_display.insert(END, "VERA: Please enter a command.\n", 'bot')
        chat_display.see(END)

def on_voice_icon_click(event=None):
    """Handle voice icon click"""
    threading.Thread(target=execute_command).start()

def set_background_image(root, image_url):
    """Set background image for the main window"""
    try:
        response = requests.get(image_url, timeout=10)
        image = Image.open(BytesIO(response.content))
        bg_image = ImageTk.PhotoImage(image)
        background_label = Label(root, image=bg_image)
        background_label.place(relwidth=1, relheight=1)
        background_label.image = bg_image
    except Exception as e:
        print(f"Failed to load background image: {e}")

def toggle_dark_mode(event=None):
    """Toggle between dark and light mode"""
    if root["bg"] == 'white':
        # Switch to dark mode
        root.configure(bg='#2e2e2e')
        header.configure(bg='#333333', fg='white')
        chat_frame.configure(bg='#3a3a3a')
        chat_display.configure(bg='#444444', fg='white')
        text_entry.configure(bg='#444444', fg='white')
        if side_panel_frame:
            side_panel_frame.configure(bg='#3a3a3a')
            for widget in side_panel_frame.winfo_children():
                if isinstance(widget, tk.Button):
                    widget.configure(bg='#555555', fg='white')
    else:
        # Switch to light mode
        root.configure(bg='white')
        header.configure(bg='black', fg='white')
        chat_frame.configure(bg='#f0f0f0')
        chat_display.configure(bg='#f8f8f8', fg='black')
        text_entry.configure(bg='white', fg='black')
        if side_panel_frame:
            side_panel_frame.configure(bg='#f0f0f0')
            for widget in side_panel_frame.winfo_children():
                if isinstance(widget, tk.Button):
                    widget.configure(bg='#007acc', fg='white')

def show_splash_screen(image_url):
    """Show splash screen on startup"""
    splash_root = tk.Tk()
    splash_root.title("Welcome")
    splash_root.geometry("1920x1080")

    try:
        response = requests.get(image_url, timeout=10)
        img_data = BytesIO(response.content)
        img = Image.open(img_data)
        photo = ImageTk.PhotoImage(img)

        background_label = Label(splash_root, image=photo)
        background_label.image = photo
        background_label.place(relwidth=1, relheight=1)

        splash_label = Label(splash_root, text="Welcome to V.E.R.A", font=("Helvetica", 30, "bold"), bg='white')
        splash_label.pack(expand=True)

        splash_root.after(2000, splash_root.destroy)
        splash_root.mainloop()
    except Exception as e:
        print(f"Could not show splash screen: {e}")
        splash_root.destroy()

def open_youtube_downloader():
    """Open YouTube downloader window"""
    if not yt_dlp:
        messagebox.showerror("Missing Library", "yt-dlp is not installed. Please install it with: pip install yt-dlp")
        return
        
    downloader_window = Toplevel(root)
    downloader_window.title("🎬 YouTube Downloader Pro")
    downloader_window.geometry("520x300")
    downloader_window.configure(bg="#eef2f5")

    if is_side_panel_open:
        toggle_side_panel()

    tk.Label(downloader_window, text="YouTube Downloader Pro", font=("Helvetica", 20, "bold"), bg="#eef2f5", fg="#333").pack(pady=10)

    url_entry = tk.Entry(downloader_window, font=("Arial", 12), width=50)
    url_entry.pack(pady=10)
    url_entry.insert(0, "Paste YouTube link here")

    format_var = tk.StringVar(value="Video (MP4)")
    format_dropdown = ttk.Combobox(downloader_window, textvariable=format_var, state="readonly",
                                   values=["Video (MP4)", "Audio (MP3)"], font=("Arial", 11), width=20)
    format_dropdown.pack(pady=5)

    status_label = tk.Label(downloader_window, text="", font=("Arial", 12), bg="#eef2f5")
    status_label.pack(pady=5)

    def start_download_internal():
        threading.Thread(target=download_video_internal, args=(url_entry, format_var, status_label)).start()

    def download_video_internal(url_entry_widget, format_var_widget, status_label_widget):
        url = url_entry_widget.get().strip()
        format_choice = format_var_widget.get()

        if not url or url == "Paste YouTube link here":
            messagebox.showwarning("Missing URL", "Please enter a YouTube URL.")
            return

        status_label_widget.config(text="Downloading...", fg="blue")

        ydl_opts = {
            'outtmpl': os.path.join(os.path.expanduser("~"), "Downloads", '%(title)s.%(ext)s'),
        }

        if format_choice == "Video (MP4)":
            ydl_opts['format'] = 'best'
        elif format_choice == "Audio (MP3)":
            ydl_opts.update({
                'format': 'bestaudio/best',
                'postprocessors': [{
                    'key': 'FFmpegExtractAudio',
                    'preferredcodec': 'mp3',
                    'preferredquality': '192',
                }]
            })

        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])
            status_label_widget.config(text="✅ Download complete!", fg="green")
        except Exception as e:
            status_label_widget.config(text="❌ Download failed.", fg="red")
            messagebox.showerror("Error", f"An error occurred:\n{e}")

    download_btn = tk.Button(downloader_window, text="⬇ Download", font=("Arial", 13, "bold"), bg="#007acc", fg="white",
                             command=start_download_internal, padx=20, pady=5)
    download_btn.pack(pady=10)

def fetch_icon(url, size=(30, 30)):
    """Fetch icon from URL"""
    try:
        response = requests.get(url, timeout=5)
        img = Image.open(BytesIO(response.content)).resize(size)
        return ImageTk.PhotoImage(img)
    except Exception as e:
        print(f"Failed to load icon from {url}: {e}")
        # Create a placeholder colored rectangle
        img = Image.new('RGB', size, color='red')
        return ImageTk.PhotoImage(img)

def create_side_panel():
    """Create side panel with navigation options"""
    global side_panel_frame, is_side_panel_open

    side_panel_frame = tk.Frame(root, bg='#f0f0f0', width=200, height=root.winfo_height())
    side_panel_frame.place(x=-200, y=0, relheight=1)

    root.bind("<Configure>", lambda e: side_panel_frame.place(x=side_panel_frame.winfo_x(), y=0, width=200, height=root.winfo_height()))

    new_chat_button = tk.Button(side_panel_frame, text="New Chat", font=("Helvetica", 12, "bold"),
                                command=new_chat, bg='#007acc', fg='white',
                                padx=10, pady=5, relief=tk.RAISED, bd=2)
    new_chat_button.pack(pady=10, fill=tk.X)

    youtube_downloader_button = tk.Button(side_panel_frame, text="YouTube Downloader", font=("Helvetica", 12, "bold"),
                                          command=open_youtube_downloader, bg='#007acc', fg='white',
                                          padx=10, pady=5, relief=tk.RAISED, bd=2)
    youtube_downloader_button.pack(pady=10, fill=tk.X)

    close_panel_button = tk.Button(side_panel_frame, text="Close Panel", font=("Helvetica", 10),
                                   command=toggle_side_panel, bg='#cc0000', fg='white',
                                   padx=5, pady=2, relief=tk.RAISED, bd=1)
    close_panel_button.pack(side=tk.BOTTOM, pady=10)

def toggle_side_panel():
    """Toggle side panel visibility"""
    global is_side_panel_open
    if is_side_panel_open:
        # Hide panel
        for i in range(0, -201, -10):
            side_panel_frame.place(x=i)
            root.update_idletasks()
            time.sleep(0.005)
        is_side_panel_open = False
    else:
        # Show panel
        for i in range(-200, 1, 10):
            side_panel_frame.place(x=i)
            root.update_idletasks()
            time.sleep(0.005)
        is_side_panel_open = True

def add_control_buttons():
    """Add control buttons to the main interface"""
    global text_entry
    icon_frame = tk.Frame(root, bg='white')
    icon_frame.pack(padx=20, pady=10)

    # Menu icon
    menu_icon = fetch_icon("https://cdn-icons-png.flaticon.com/512/545/545667.png")
    menu_label = Label(icon_frame, image=menu_icon, bg='white', cursor='hand2')
    menu_label.image = menu_icon
    menu_label.pack(side=tk.LEFT, padx=5)
    menu_label.bind("<Button-1>", lambda e: toggle_side_panel())

    # Text entry
    text_entry = Entry(icon_frame, font=("Helvetica", 14), bg='white', width=50)
    text_entry.insert(0, "Type your command here...")
    text_entry.bind("<FocusIn>", lambda args: text_entry.delete('0', 'end'))
    text_entry.pack(side=tk.LEFT, padx=10)

    # Control icons
    mic_icon = fetch_icon("https://icones.pro/wp-content/uploads/2021/12/icone-de-microphone-bleue.png")
    stop_icon = fetch_icon("https://cdn-icons-png.flaticon.com/512/4029/4029084.png")
    add_file_icon = fetch_icon("https://cdn-icons-png.flaticon.com/512/3097/3097412.png")

    # Add file button
    add_file_label = Label(icon_frame, image=add_file_icon, bg='white', cursor='hand2')
    add_file_label.image = add_file_icon
    add_file_label.pack(side=tk.RIGHT, padx=10)
    add_file_label.bind("<Button-1>", lambda e: upload_file())

    # Stop button
    stop_label = Label(icon_frame, image=stop_icon, bg='white', cursor='hand2')
    stop_label.image = stop_icon
    stop_label.pack(side=tk.RIGHT, padx=10)
    stop_label.bind("<Button-1>", lambda e: stop_all())

    # Microphone button
    mic_label = Label(icon_frame, image=mic_icon, bg='white', cursor='hand2')
    mic_label.image = mic_icon
    mic_label.pack(side=tk.RIGHT, padx=10)
    mic_label.bind("<Button-1>", on_voice_icon_click)

# --- NEW: Function to handle application closing ---
def on_closing():
    """Gracefully closes the application and cancels any running threads."""
    global timeout_thread
    if timeout_thread:
        timeout_thread.cancel()
    root.destroy()

# Global variable for last response
last_response = ""

# Main application setup
root = tk.Tk()
root.title("V.E.R.A")
root.geometry("700x700")
root.configure(bg='white')

# Set background image
set_background_image(root, 'https://wallpapercave.com/wp/wp1913256.jpg')
# Header
header = Label(root, text="V.E.R.A", font=("Helvetica", 25, "bold"), bg='black', fg='white', pady=10)
header.pack(fill=tk.X)

# Chat frame
chat_frame = tk.Frame(root, bg='#f0f0f0')
chat_frame.pack(pady=10)

# Scrollbar
scrollbar = Scrollbar(chat_frame)
scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

# Chat display
chat_display = Text(chat_frame, yscrollcommand=scrollbar.set, wrap=tk.WORD, bg='#f8f8f8', fg='black', 
                   font=("Helvetica", 14), bd=0, padx=10, pady=10)
chat_display.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
scrollbar.config(command=chat_display.yview)

# Configure text tags
chat_display.tag_config('user', foreground='#1e90ff', font=("Helvetica", 14, 'bold'))
chat_display.tag_config('bot', foreground='#32cd32', font=("Helvetica", 14, 'italic'))
chat_display.tag_config('timestamp', foreground='#a9a9a9', font=("Helvetica", 10, 'italic'))

# Add control buttons
add_control_buttons()

# Create side panel
create_side_panel()

# Initialize
greet_user()
load_history()

# Bind keyboard shortcuts
root.bind("<Control-d>", toggle_dark_mode)
root.bind("<Return>", lambda event=None: on_text_command())
# --- MODIFIED: Handle window close event gracefully ---
root.protocol("WM_DELETE_WINDOW", on_closing)

# --- NEW: Start the background listener thread for wake word ---
listener_thread = threading.Thread(target=listen_for_wake_word, daemon=True)
listener_thread.start()

# Start main loop
root.mainloop()
